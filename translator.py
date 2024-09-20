import re
from pptx import Presentation
from pptx.util import Pt
from deep_translator import GoogleTranslator

class PPTTranslator:
    def __init__(self, file_path, save_path):
        self.file_path = file_path
        self.save_path = save_path
        self.translator = GoogleTranslator(source='zh-CN', target='ko')

    def translate(self):
        try:
            prs = Presentation(self.file_path)
            total_slides = len(prs.slides)
            for index, slide in enumerate(prs.slides):
                self.process_slide(slide)
                # Calculate progress (assumes 0-100% range)
                yield int((index + 1) / total_slides * 100)
            prs.save(self.save_path)
            return None  # No error
        except Exception as e:
            return f"翻译错误: {e}"

    def process_slide(self, slide):
        for shape in slide.shapes:
            if shape.has_text_frame:
                self.process_text_frame(shape.text_frame, shape)

    def process_text_frame(self, text_frame, shape):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if not self.contains_numbers(run.text):
                    translated_text = self.translator.translate(run.text)
                    run.text = f"{run.text}\n{translated_text}"
                self.adjust_font_size(run)
        self.fit_text_in_shape(text_frame, shape)

    def contains_numbers(self, text):
        return bool(re.search(r'\d+', text))

    def adjust_font_size(self, run):
        if run.font.size:
            original_size = run.font.size.pt
            run.font.size = Pt(original_size)

    def fit_text_in_shape(self, text_frame, shape):
        max_width = shape.width
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(14)
        self.adjust_text_to_fit(text_frame, max_width)

    def adjust_text_to_fit(self, text_frame, max_width):
        for paragraph in text_frame.paragraphs:
            while self.get_text_width(paragraph) > max_width:
                for run in paragraph.runs:
                    if run.font.size:
                        run.font.size = Pt(run.font.size.pt - 1)
                        if run.font.size.pt < 6:
                            break
                if all(run.font.size.pt < 6 for run in paragraph.runs):
                    break

    def get_text_width(self, paragraph):
        return len(paragraph.text) * 6
